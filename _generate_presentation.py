# Generates the Lexora project presentation (.pptx), dark theme.
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# ---- palette (matches Lexora's dark UI) --------------------------------------
BG = RGBColor(0x12, 0x14, 0x1A)          # near-black background
PANEL = RGBColor(0x1C, 0x1F, 0x2A)       # card background
ACCENT = RGBColor(0x6C, 0x8C, 0xFF)      # soft indigo
ACCENT2 = RGBColor(0x3D, 0xD6, 0x8C)     # green (offline / success)
WARN = RGBColor(0xFF, 0x5C, 0x5C)        # red (contradictions)
TEXT = RGBColor(0xEA, 0xEC, 0xF2)        # near-white
MUTED = RGBColor(0x9A, 0xA3, 0xB5)       # grey text
FONT = "Segoe UI"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height


def add_slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def box(slide, x, y, w, h, fill=None, line=None):
    sh = slide.shapes.add_shape(1, x, y, w, h)  # 1 = rectangle
    sh.fill.solid() if fill else sh.fill.background()
    if fill:
        sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def rounded(slide, x, y, w, h, fill=PANEL):
    sh = slide.shapes.add_shape(5, x, y, w, h)  # 5 = rounded rectangle
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    try:
        sh.adjustments[0] = 0.08
    except Exception:
        pass
    return sh


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is a list of (text, size, color, bold)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para_runs in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        if line_spacing:
            p.line_spacing = line_spacing
        for t, size, color, bold in para_runs:
            r = p.add_run()
            r.text = t
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
    return tb


def title_bar(slide, title, subtitle=None):
    box(slide, Inches(0.55), Inches(0.62), Inches(0.09), Inches(0.62), fill=ACCENT)
    runs = [[(title, 30, TEXT, True)]]
    if subtitle:
        runs.append([(subtitle, 14, MUTED, False)])
    text(slide, Inches(0.85), Inches(0.42), Inches(11.9), Inches(1.2), runs, space_after=2)


def footer(slide, n):
    text(slide, Inches(0.55), Inches(7.02), Inches(6), Inches(0.4),
         [[("LEXORA", 10, ACCENT, True), ("  \u2022  Offline-First AI Assistant for Lawyers",
                                          10, MUTED, False)]])
    text(slide, Inches(12.3), Inches(7.02), Inches(0.7), Inches(0.4),
         [[(str(n), 10, MUTED, False)]], align=PP_ALIGN.RIGHT)


def bullets(slide, x, y, w, h, items, size=15, gap=8):
    """items: list of (head, body) or plain strings."""
    runs = []
    for it in items:
        if isinstance(it, tuple):
            head, body = it
            runs.append([("\u25AA  ", size, ACCENT, True), (head + " \u2014 ", size, TEXT, True),
                         (body, size, MUTED, False)])
        else:
            runs.append([("\u25AA  ", size, ACCENT, True), (it, size, TEXT, False)])
    text(slide, x, y, w, h, runs, space_after=gap, line_spacing=1.05)


def card(slide, x, y, w, h, head, body, head_color=ACCENT, body_size=12.5, head_size=15):
    rounded(slide, x, y, w, h)
    text(slide, x + Inches(0.22), y + Inches(0.14), w - Inches(0.44), h - Inches(0.28),
         [[(head, head_size, head_color, True)], [(body, body_size, MUTED, False)]],
         space_after=4, line_spacing=1.05)


# ============================ 1. TITLE ========================================
s = add_slide()
box(s, 0, 0, W, Emu(int(H * 0.012)), fill=ACCENT)
text(s, Inches(1), Inches(2.0), Inches(11.3), Inches(1.8),
     [[("LEXORA", 66, TEXT, True)],
      [("An Offline-First AI Assistant for Lawyers", 26, ACCENT, False)]],
     align=PP_ALIGN.CENTER, space_after=8)
text(s, Inches(1), Inches(4.35), Inches(11.3), Inches(1.4),
     [[("Private  \u2022  On-Device LLM  \u2022  Zero Data Leaves the Machine", 16, ACCENT2, True)],
      [("", 8, MUTED, False)],
      [("Taabish Ahmed Ansari", 18, TEXT, True)],
      [("Faculty of Science & Technology, IFHE University  \u2022  June\u2013July 2026",
        13, MUTED, False)]],
     align=PP_ALIGN.CENTER, space_after=4)

# ============================ 2. THE PROBLEM ==================================
s = add_slide()
title_bar(s, "The Problem", "Why lawyers can't just use ChatGPT")
bullets(s, Inches(0.85), Inches(1.65), Inches(11.6), Inches(4.0), [
    ("Confidentiality", "attorney\u2013client privilege makes it unacceptable to send case "
     "data to cloud AI services \u2014 every popular assistant is cloud-hosted."),
    ("Fragmented memory", "client facts live across notebooks, e-mails, and memory; nothing "
     "remembers everything a client ever said and recalls it instantly."),
    ("Missed contradictions", "clients contradict themselves weeks apart (\u201cI was at the "
     "office at 5 PM\u201d vs \u201cI was not at the office\u201d) \u2014 catching this "
     "manually across long case histories is error-prone."),
    ("Administrative overhead", "summarizing calls, extracting action items, scheduling "
     "hearings, and e-mailing summaries eat into billable legal work."),
], size=16, gap=14)
rounded(s, Inches(0.85), Inches(5.85), Inches(11.6), Inches(0.85), fill=PANEL)
text(s, Inches(1.1), Inches(5.98), Inches(11.1), Inches(0.65),
     [[("The gap:  ", 15, WARN, True),
       ("lawyers need AI-grade assistance with a guarantee that no client data ever leaves "
        "their machine.", 15, TEXT, False)]], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 2)

# ============================ 3. THE SOLUTION =================================
s = add_slide()
title_bar(s, "The Solution: Lexora", "Everything that touches client data runs locally")
cw, ch = Inches(3.7), Inches(1.55)
xs = [Inches(0.85), Inches(4.8), Inches(8.75)]
ys = [Inches(1.7), Inches(3.45), Inches(5.2)]
cards = [
    ("Offline Chat", "ChatGPT-style legal assistant powered by TinyLlama-1.1B running "
     "on-device (GPU fp16 or CPU)."),
    ("Persistent Memory", "Every client statement is embedded and indexed \u2014 recalled "
     "semantically in future chats, survives restarts."),
    ("Knowledge Graph", "Entities & relations auto-extracted per client, rendered with a "
     "custom Canvas force-directed view."),
    ("Offline RAG", "Upload PDF case files; answers cite sources. Local embeddings + custom "
     "NumPy vector store."),
    ("Contradiction Detection", "Deterministic engine flags polarity and time/number "
     "conflicts with prior testimony \u2014 red warning card."),
    ("AI Call Assistant", "Live speech-to-text in the browser, real-time memory commits, "
     "auto meeting summary + action items."),
    ("Smart Scheduling", "Natural-language dates (\u201chearing next Friday 5pm\u201d) \u2192 "
     "conflict check + one-click Google Calendar link."),
    ("Case Dashboard", "Clients by status & type, totals for calls, statements, documents, "
     "contradictions \u2014 with an activity audit trail."),
    ("Client Management", "Create, update, delete clients; deletion purges chat, graph, "
     "calls, and every vector everywhere."),
]
for i, (head, body) in enumerate(cards):
    hc = ACCENT2 if head in ("Offline Chat", "Offline RAG") else (
        WARN if head == "Contradiction Detection" else ACCENT)
    card(s, xs[i % 3], ys[i // 3], cw, ch, head, body, head_color=hc, body_size=11.5,
         head_size=14)
footer(s, 3)

# ============================ 4. ARCHITECTURE =================================
s = add_slide()
title_bar(s, "System Architecture", "A clean layered design \u2014 JSON on disk, no external services")
layers = [
    ("Frontend  \u2014  static/", "Vanilla HTML / CSS / JS  \u2022  dark ChatGPT-style UI  "
     "\u2022  HTML5 Canvas graph  \u2022  Chart.js (vendored)  \u2022  Web Speech API", ACCENT),
    ("API  \u2014  app.py", "FastAPI + Uvicorn + Pydantic  \u2022  REST endpoints: chat, "
     "clients, upload, live, call, dashboard  \u2022  Twilio webhooks  \u2022  no-cache "
     "middleware", ACCENT),
    ("Orchestration  \u2014  chat.py", "Builds each turn: memory recall + RAG chunks + recent "
     "turns + contradiction \u2192 LLM \u2192 clean \u2192 graph update \u2192 persist", ACCENT2),
    ("Core engines  \u2014  llm.py  \u2022  rag.py  \u2022  memory.py",
     "TinyLlama generation + JSON extraction  \u2022  BGE embeddings + NumPy cosine index + "
     "pypdf ingest  \u2022  clients, graph, contradiction engine, activity log", ACCENT2),
    ("Persistence  \u2014  data/  &  chroma_db/", "store.json (all case data)  \u2022  "
     "documents.json + memory.json (vectors)  \u2022  thread-safe (RLock), flushed on every "
     "write", MUTED),
]
y = Inches(1.62)
for head, body, hc in layers:
    rounded(s, Inches(0.85), y, Inches(11.6), Inches(0.92))
    text(s, Inches(1.1), y + Inches(0.08), Inches(11.1), Inches(0.8),
         [[(head, 14, hc, True)], [(body, 11.5, MUTED, False)]], space_after=2)
    y += Inches(1.04)
text(s, Inches(0.85), y + Inches(0.02), Inches(11.6), Inches(0.4),
     [[("Side modules: ", 12, TEXT, True),
       ("calendarx.py (NL scheduling)  \u2022  mailer.py (Gmail SMTP)  \u2022  telephony.py "
        "(Twilio real calls)  \u2022  seed_data.py (demo data + PDF generator)", 12, MUTED,
        False)]])
footer(s, 4)

# ============================ 5. CHAT TURN FLOW ===============================
s = add_slide()
title_bar(s, "Anatomy of a Chat Turn", "What happens when the lawyer sends one message")
steps = [
    ("1. Store", "Message saved as a client statement, embedded, and added to the memory "
     "vector index."),
    ("2. Check", "Contradiction engine recalls similar past statements and runs polarity + "
     "numeric/time conflict rules."),
    ("3. Assemble", "Context = semantic memory recall + top-k RAG chunks + recent turns + "
     "any detected contradiction."),
    ("4. Generate", "TinyLlama produces the reply (chat template, temperature / top-p / "
     "repetition penalty); regex strips hallucinated URLs and context echo."),
    ("5. Learn", "Entities & relations extracted as JSON, sanitized, merged "
     "case-insensitively into the client's knowledge graph."),
    ("6. Persist", "Everything flushed to store.json and the JSON vector files \u2014 close "
     "the app, reopen, nothing is lost."),
]
y = Inches(1.65)
for i, (head, body) in enumerate(steps):
    rounded(s, Inches(0.85), y, Inches(11.6), Inches(0.74))
    text(s, Inches(1.1), y + Inches(0.09), Inches(2.0), Inches(0.55),
         [[(head, 15, ACCENT, True)]])
    text(s, Inches(3.2), y + Inches(0.09), Inches(9.0), Inches(0.55),
         [[(body, 12.5, TEXT, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.86)
footer(s, 5)

# ============================ 6. LLM LAYER ====================================
s = add_slide()
title_bar(s, "The Local LLM Layer", "Small model, big guardrails")
bullets(s, Inches(0.85), Inches(1.65), Inches(6.6), Inches(4.6), [
    ("TinyLlama-1.1B-Chat", "loaded lazily via HuggingFace Transformers + PyTorch; fp16 on "
     "CUDA GPU, fp32 on CPU."),
    ("Native chat template", "with sampling controls \u2014 temperature, top-p, repetition "
     "penalty."),
    ("Output hygiene", "regex post-processing strips hallucinated URLs and context-echo "
     "lines small models emit."),
    ("Structured generation", "generate_json() extracts the first balanced JSON object from "
     "raw output \u2014 tolerates code fences and prose."),
    ("Silent fallback", "optional, env-gated Gemini fallback only on local failure; never "
     "shown in the UI."),
], size=14, gap=10)
rounded(s, Inches(7.75), Inches(1.65), Inches(4.7), Inches(4.6), fill=PANEL)
text(s, Inches(8.0), Inches(1.85), Inches(4.2), Inches(4.2),
     [[("Design principle", 16, ACCENT2, True)],
      [("", 6, MUTED, False)],
      [("Treat the model as one component in a pipeline \u2014 not an oracle.", 15, TEXT, True)],
      [("", 6, MUTED, False)],
      [("Deterministic engines back the 1.1B model on every must-be-correct task: "
        "contradictions, action items, graph labels.", 13, MUTED, False)],
      [("", 6, MUTED, False)],
      [("Result: demos stay reliable regardless of the small model's variance.", 13, MUTED,
        False)]], space_after=4)
footer(s, 6)

# ============================ 7. RAG ==========================================
s = add_slide()
title_bar(s, "Offline RAG over Case Documents", "PDF in \u2192 cited answers out \u2014 no internet")
steps = ["PDF upload", "pypdf text extract", "Chunk 900/150 overlap", "BGE embeddings",
         "NumPy cosine index", "Top-k cited answer"]
x = Inches(0.85)
bw = Inches(1.78)
for i, st in enumerate(steps):
    rounded(s, x, Inches(1.75), bw, Inches(0.95), fill=PANEL)
    text(s, x + Inches(0.08), Inches(1.78), bw - Inches(0.16), Inches(0.9),
         [[(st, 12, TEXT, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < len(steps) - 1:
        text(s, x + bw - Inches(0.06), Inches(1.86), Inches(0.35), Inches(0.7),
             [[("\u203A", 22, ACCENT, True)]], anchor=MSO_ANCHOR.MIDDLE)
    x += bw + Inches(0.18)
rounded(s, Inches(0.85), Inches(3.1), Inches(11.6), Inches(1.7), fill=PANEL)
text(s, Inches(1.1), Inches(3.28), Inches(11.1), Inches(1.4),
     [[("War story: why a custom vector store?", 15, WARN, True)],
      [("The installed ChromaDB build threw a fatal mismatched-types (u64/BLOB) compaction "
        "error on a fresh database \u2014 silently dropping every vector. Replaced it with a "
        "dependency-light NumPy cosine index (vectors normalized \u2192 cosine = dot product) "
        "persisted as JSON. Same retrieval quality, zero failure points.", 12.5, MUTED,
        False)]], space_after=5)
bullets(s, Inches(0.85), Inches(5.05), Inches(11.6), Inches(1.6), [
    ("Per-client filtering", "retrieval is scoped to the active client, falling back to the "
     "shared library."),
    ("Source citation", "every answer names the document it came from."),
], size=14, gap=8)
footer(s, 7)

# ============================ 8. MEMORY + GRAPH ===============================
s = add_slide()
title_bar(s, "Memory & Knowledge Graph", "Lexora remembers \u2014 and connects the dots")
bullets(s, Inches(0.85), Inches(1.65), Inches(6.6), Inches(4.8), [
    ("Semantic recall", "each statement is embedded; future chats retrieve the most similar "
     "past statements into context \u2014 ChatGPT-like memory that survives restarts."),
    ("Auto-extraction", "TinyLlama returns {entities, relations} JSON per factual message."),
    ("Sanitization", "heuristics drop hallucinated URLs and placeholder labels "
     "(\u201cPERSON\u201d, \u201cDATE\u201d)."),
    ("Smart merging", "case-insensitive node/edge merge \u2014 the same entity across "
     "different conversations connects automatically; orphans attach to the client node."),
    ("No graph DB", "plain {nodes, edges} lists in store.json."),
], size=14, gap=10)
rounded(s, Inches(7.75), Inches(1.65), Inches(4.7), Inches(4.8), fill=PANEL)
text(s, Inches(8.0), Inches(1.85), Inches(4.2), Inches(4.4),
     [[("Custom Canvas renderer", 16, ACCENT, True)],
      [("", 5, MUTED, False)],
      [("A force-directed graph built from scratch on HTML5 Canvas:", 13, TEXT, False)],
      [("", 4, MUTED, False)],
      [("\u25AA  node\u2013node repulsion", 13, MUTED, False)],
      [("\u25AA  spring attraction along edges", 13, MUTED, False)],
      [("\u25AA  center gravity", 13, MUTED, False)],
      [("\u25AA  drag-to-move interaction", 13, MUTED, False)],
      [("\u25AA  color-coded entity types", 13, MUTED, False)],
      [("", 5, MUTED, False)],
      [("Zero external libraries \u2014 no D3, no vis.js, no three.js.", 13, ACCENT2, True)]],
     space_after=3)
footer(s, 8)

# ============================ 9. CONTRADICTIONS ===============================
s = add_slide()
title_bar(s, "Contradiction Detection", "Deterministic by design \u2014 too important to leave to a 1.1B model")
rounded(s, Inches(0.85), Inches(1.65), Inches(5.6), Inches(2.3), fill=PANEL)
text(s, Inches(1.1), Inches(1.82), Inches(5.1), Inches(2.0),
     [[("Rule 1 \u2014 Polarity flip", 15, WARN, True)],
      [("Same subject (shared tokens + cosine \u2265 0.45), one statement affirms, the other "
        "denies.", 12.5, MUTED, False)],
      [("", 3, MUTED, False)],
      [("\u201cI was at the office\u201d  vs  \u201cI was NOT at the office\u201d", 12.5,
        TEXT, True)]], space_after=5)
rounded(s, Inches(6.85), Inches(1.65), Inches(5.6), Inches(2.3), fill=PANEL)
text(s, Inches(7.1), Inches(1.82), Inches(5.1), Inches(2.0),
     [[("Rule 2 \u2014 Number / time conflict", 15, WARN, True)],
      [("Same subject but different numbers or clock times, while skipping corroborating "
        "evidence.", 12.5, MUTED, False)],
      [("", 3, MUTED, False)],
      [("\u201cleft at 5 PM\u201d  vs  \u201cleft at 9 PM\u201d", 12.5, TEXT, True)]],
     space_after=5)
bullets(s, Inches(0.85), Inches(4.35), Inches(11.6), Inches(2.2), [
    ("Pipeline", "new statement \u2192 embedded \u2192 most similar prior statements recalled "
     "from the client's memory index \u2192 rules applied."),
    ("Explanation", "the chatbot explains the conflict in natural language (TinyLlama)."),
    ("Guaranteed visual", "a red warning card always appears in the UI when a contradiction "
     "is found."),
], size=14, gap=10)
footer(s, 9)

# ============================ 10. CALL ASSISTANT ==============================
s = add_slide()
title_bar(s, "AI Call Assistant & Live Calls", "The assistant listens, remembers, and writes the minutes")
bullets(s, Inches(0.85), Inches(1.65), Inches(11.6), Inches(4.9), [
    ("Private speech-to-text", "browser-native Web Speech API (Chrome/Edge) \u2014 no audio "
     "ever leaves the machine; for a real two-phone call, the phone goes on speaker next to "
     "the laptop."),
    ("Real-time pipeline", "each finalized utterance hits /api/live \u2192 committed to "
     "memory + knowledge graph instantly, with live contradiction checks."),
    ("Auto client creation", "a name detector (\u201cmy name is \u2026\u201d) plus a "
     "keyword case-type classifier (theft \u2192 Criminal, divorce \u2192 Divorce) creates a "
     "new client mid-call when none is selected."),
    ("End of call", "TinyLlama writes the meeting summary; action items extracted via LLM "
     "JSON with a robust regex fallback (\u201cwill\u201d, \u201csend by Friday\u201d, "
     "\u201cfollow up\u201d)."),
    ("Real telephony (optional)", "Twilio integration with server-side transcription "
     "callbacks for genuine phone calls."),
], size=14.5, gap=13)
footer(s, 10)

# ============================ 11. PRODUCTIVITY EXTRAS =========================
s = add_slide()
title_bar(s, "Scheduling, E-mail & Dashboard", "Rounding out the product")
card(s, Inches(0.85), Inches(1.7), Inches(3.75), Inches(4.6), "Smart Scheduling",
     "Detects scheduling intent via cue words \u2014 and excludes past testimony "
     "(\u201che said he met her Friday\u201d is not an appointment).\n\nRegex pulls clock "
     "times first so \u201c11am\u201d isn't parsed as day 11; dateparser resolves the date "
     "(future-preferring).\n\nConflict check + one-click \u201cAdd to Google Calendar\u201d "
     "link.", body_size=12)
card(s, Inches(4.8), Inches(1.7), Inches(3.75), Inches(4.6), "Case Summary E-mails",
     "One click sends a case summary to the client via Gmail SMTP (app password).\n\n"
     "Falls back to a default recipient for demo clients with no e-mail on file.",
     head_color=ACCENT2, body_size=12)
card(s, Inches(8.75), Inches(1.7), Inches(3.75), Inches(4.6), "Case Insights Dashboard",
     "Clients aggregated by status (solved / unsolved / \u2026) and case type.\n\nTotals for "
     "calls, statements, documents, contradictions.\n\nChart.js doughnut + bar charts on a "
     "dark theme.\n\nActivity log: who-did-what-when for every action.", body_size=12)
footer(s, 11)

# ============================ 12. TECH STACK ==================================
s = add_slide()
title_bar(s, "Technology Stack", "Deliberately lightweight \u2014 everything local, everything inspectable")
rows = [
    ("Language / Runtime", "Python 3.11 (CUDA-enabled)"),
    ("LLM", "TinyLlama-1.1B-Chat \u2014 HuggingFace Transformers + PyTorch (fp16 GPU / CPU)"),
    ("Embeddings", "BGE-small-en-v1.5 via sentence-transformers \u2014 384-dim, normalized"),
    ("Vector store", "Custom NumPy cosine index, persisted as JSON (replaced ChromaDB)"),
    ("Backend", "FastAPI + Uvicorn + Pydantic  \u2022  pypdf  \u2022  python-multipart"),
    ("Frontend", "Vanilla HTML/CSS/JS  \u2022  HTML5 Canvas  \u2022  Chart.js (vendored)  "
     "\u2022  Web Speech API"),
    ("Integrations", "Gmail SMTP  \u2022  Google Calendar links  \u2022  Twilio voice "
     "(optional)  \u2022  dateparser"),
    ("Persistence", "JSON files \u2014 store.json + documents.json + memory.json "
     "(human-readable, portable)"),
]
y = Inches(1.62)
for head, body in rows:
    rounded(s, Inches(0.85), y, Inches(11.6), Inches(0.56), fill=PANEL)
    text(s, Inches(1.1), y + Inches(0.05), Inches(2.9), Inches(0.45),
         [[(head, 13, ACCENT, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(4.1), y + Inches(0.05), Inches(8.1), Inches(0.45),
         [[(body, 12.5, TEXT, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.65)
footer(s, 12)

# ============================ 13. CHALLENGES ==================================
s = add_slide()
title_bar(s, "Challenges & Learnings", "What building Lexora taught me")
bullets(s, Inches(0.85), Inches(1.65), Inches(11.6), Inches(4.9), [
    ("Taming a 1.1B model", "small LLMs hallucinate URLs, echo context, emit broken JSON. "
     "Learning: wrap the model in layered defenses \u2014 output hygiene, balanced-JSON "
     "extraction, label sanitization, deterministic engines for correctness-critical tasks."),
    ("The ChromaDB failure", "a fatal compaction bug silently dropped every vector. "
     "Learning: understanding the math (normalized embeddings \u2192 cosine = dot product) "
     "made the heavyweight dependency unnecessary \u2014 40 lines of NumPy replaced it."),
    ("Natural-language dates", "\u201cmeet at 11am\u201d was parsed as \u201cday 11\u201d. "
     "Learning: extract clock times with regex before date parsing; prefer future dates; "
     "exclude testimony sentences from scheduling."),
    ("Data hygiene", "deleting a client must purge chat, statements, graph, calls, and "
     "vectors in two indexes. Learning: design deletion paths as carefully as creation "
     "paths."),
], size=14, gap=13)
footer(s, 13)

# ============================ 14. CONCLUSION ==================================
s = add_slide()
title_bar(s, "Conclusion & Future Work")
rounded(s, Inches(0.85), Inches(1.65), Inches(11.6), Inches(1.5), fill=PANEL)
text(s, Inches(1.1), Inches(1.82), Inches(11.1), Inches(1.2),
     [[("Lexora proves a genuinely useful, private AI assistant for lawyers can run entirely "
        "on local hardware \u2014 ", 15, TEXT, False),
       ("conversational help, persistent memory, document intelligence, contradiction "
        "detection, call summaries, scheduling, and analytics, with zero bytes of client "
        "data leaving the machine.", 15, ACCENT2, True)]], space_after=4)
text(s, Inches(0.85), Inches(3.45), Inches(11.6), Inches(0.5),
     [[("Future work", 18, ACCENT, True)]])
bullets(s, Inches(0.85), Inches(4.0), Inches(11.6), Inches(2.5), [
    "Upgrade to newer small models (Phi-3, Llama-3.2-1B) and fine-tune on legal corpora.",
    "OCR for scanned filings and richer document understanding.",
    "Multi-user support with authentication and role-based access.",
    "Package as a signed desktop installer for non-technical users.",
], size=14.5, gap=9)
footer(s, 14)

# ============================ 15. THANK YOU ===================================
s = add_slide()
box(s, 0, 0, W, Emu(int(H * 0.012)), fill=ACCENT)
text(s, Inches(1), Inches(2.5), Inches(11.3), Inches(1.4),
     [[("Thank You", 54, TEXT, True)],
      [("Questions & Demo", 20, ACCENT, False)]], align=PP_ALIGN.CENTER, space_after=10)
text(s, Inches(1), Inches(4.6), Inches(11.3), Inches(1.0),
     [[("Taabish Ahmed Ansari", 17, TEXT, True)],
      [("LEXORA \u2014 An Offline-First AI Assistant for Lawyers", 13, MUTED, False)]],
     align=PP_ALIGN.CENTER, space_after=4)

OUT = "Lexora Presentation - Taabish Ahmed Ansari.pptx"
prs.save(OUT)
print("saved:", OUT, "| slides:", len(prs.slides.__iter__.__self__._sldIdLst))
